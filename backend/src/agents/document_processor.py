from openai import OpenAI
import json
import io
import re
import logging
from src.config import settings
from src.config import prompts
from src.utils.ai_retry import chat_with_retry

logger = logging.getLogger("MudarisDocProcessor")


class DocumentProcessorAgent:
    def __init__(self):
        self.client = OpenAI(
            api_key=settings.OPENROUTER_API_KEY,
            base_url=settings.OPENROUTER_BASE_URL,
            default_headers={
                "HTTP-Referer": "https://mudaris-app.com",
                "X-Title": "Mudaris AI Engine"
            }
        )
        self.model_id = settings.OPENROUTER_MODEL

    def extract_text(self, file_bytes: bytes, file_type: str) -> str:
        if file_type == "pdf":
            return self._extract_pdf(file_bytes)
        elif file_type == "pptx":
            return self._extract_pptx(file_bytes)
        elif file_type == "docx":
            return self._extract_docx(file_bytes)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _extract_pdf(self, file_bytes: bytes) -> str:
        # Robust extraction: PyMuPDF first (handles tricky font encodings that
        # PyPDF2 returns empty for), with PyPDF2 as a fallback.
        from src.utils.pdf_extract import extract_pdf_text
        return extract_pdf_text(file_bytes)

    def _extract_pptx(self, file_bytes: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    def _extract_docx(self, file_bytes: bytes) -> str:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        return "\n\n".join(paragraphs)

    def analyze_document(self, text: str, course_title: str, image_uris: list = None) -> dict:
        """Run AI analysis on document text. Retries once on malformed JSON,
        and RAISES on persistent failure (so the upload handler can mark the
        document 'failed' instead of silently storing an empty analysis).
        When page images are supplied, the (vision) model also reads equations,
        diagrams, figures, and code that plain text extraction misses.
        """
        from src.utils.json_parse import parse_with_retry

        logger.info(f"Analyzing document for course: {course_title}"
                    + (f" (+{len(image_uris)} page images)" if image_uris else ""))
        truncated = text[:15000]

        def _user_content():
            if image_uris:
                parts = [{
                    "type": "text",
                    "text": (
                        f"Course: {course_title}\n\nThe lecture pages are shown as images below"
                        " (read the equations, diagrams, figures, and code in them)."
                        + (f"\n\n---EXTRACTED TEXT (for reference)---\n{truncated}" if truncated.strip() else "")
                    ),
                }]
                for uri in image_uris[:8]:
                    parts.append({"type": "image_url", "image_url": {"url": uri}})
                return parts
            return f"Course: {course_title}\n\n---DOCUMENT TEXT---\n{truncated}"

        def _call(reminder: str) -> str:
            response = chat_with_retry(
                self.client,
                model=self.model_id,
                messages=[
                    {
                        "role": "system",
                        "content": prompts.DOCUMENT_PROCESSOR_SYSTEM_PROMPT + reminder,
                    },
                    {"role": "user", "content": _user_content()},
                ],
                temperature=0.1,
                max_tokens=settings.OPENROUTER_MAX_TOKENS,
            )
            return (response.choices[0].message.content or "").strip()

        result = parse_with_retry(
            _call,
            required_keys=(
                "topics",
                "definitions",
                "formulas",
                "chapterMapping",
                "keyConceptCount",
            ),
            label="document analysis",
        )

        # Sanity check: if every output field is empty, the analysis was
        # effectively useless — fail loudly so the user can retry.
        topics = result.get("topics") or []
        defs = result.get("definitions") or []
        formulas = result.get("formulas") or []
        chapters = result.get("chapterMapping") or []
        if not topics and not defs and not formulas and not chapters:
            raise RuntimeError(
                "The AI returned an empty analysis (no topics, definitions, "
                "formulas, or chapters). Try re-analyzing — the model may "
                "have been overloaded or had trouble with this content."
            )

        return result

    def analyze_tutorial(self, text: str, course_title: str, image_uris: list = None) -> dict:
        """Analyze a TUTORIAL / practice-problem sheet. Unlike a lecture document,
        a tutorial is a set of QUESTIONS the student solves — so we capture each
        problem (what it asks, the given data/tables, the concept, and the method)
        precisely enough that an exam can regenerate a variant with new numbers.
        Page images let the model read tables, equations, diagrams, and code."""
        from src.utils.json_parse import parse_with_retry

        logger.info(f"Analyzing tutorial for course: {course_title}"
                    + (f" (+{len(image_uris)} page images)" if image_uris else ""))
        truncated = text[:15000]

        def _user_content():
            if image_uris:
                parts = [{
                    "type": "text",
                    "text": (
                        f"Course: {course_title}\n\nThe tutorial pages are shown as images below"
                        " (read the problems, tables, equations, diagrams, and code in them)."
                        + (f"\n\n---EXTRACTED TEXT (for reference)---\n{truncated}" if truncated.strip() else "")
                    ),
                }]
                for uri in image_uris[:8]:
                    parts.append({"type": "image_url", "image_url": {"url": uri}})
                return parts
            return f"Course: {course_title}\n\n---TUTORIAL TEXT---\n{truncated}"

        def _call(reminder: str) -> str:
            response = chat_with_retry(
                self.client,
                model=self.model_id,
                messages=[
                    {
                        "role": "system",
                        "content": prompts.TUTORIAL_SYSTEM_PROMPT + reminder,
                    },
                    {"role": "user", "content": _user_content()},
                ],
                temperature=0.1,
                # Tutorials can have many problems + code; give the JSON room so
                # it isn't truncated mid-output.
                max_tokens=max(settings.OPENROUTER_MAX_TOKENS, 8000),
            )
            return (response.choices[0].message.content or "").strip()

        result = parse_with_retry(
            _call,
            required_keys=("topics", "problems", "formulas", "skills", "problemCount"),
            label="tutorial analysis",
        )

        problems = result.get("problems") or []
        topics = result.get("topics") or []
        if not problems and not topics:
            raise RuntimeError(
                "The AI returned an empty tutorial analysis (no problems or "
                "topics). Try re-analyzing — the model may have been overloaded."
            )

        return result
